# -*- coding: utf-8 -*-
"""
Created on Thu Jan 23 16:49:45 2020
@author: Bharathiraja_A
"""
from nltk.tokenize import word_tokenize
from pptx import Presentation
from docx import Document
import os
import re

input_words= []
abbr_list = []
text_runs = []

input_file = 'file.docx'
extn = os.path.splitext(input_file)[-1].lower()

if extn == ".docx":
    document = Document(input_file)
    for para in document.paragraphs:
        words = word_tokenize(para.text)
        input_words=input_words+words
        generate=1
elif extn == ".pptx":
    prs = Presentation(input_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    for val in text_runs:
        words = word_tokenize(val)
        input_words=input_words+words
        generate=1
else:
    print("unsupported format")
    generate=0

def acronymsearch(text):
    patterns = '^[A-Z]{2,}$'
    if re.search(patterns,  text):
        return text
    else:
        return 0;
if generate==1:
    for x in input_words:
        res = acronymsearch(x)
        if res!=0:
            abbr_list.append(res)
    
    #Removes duplicates by converting into Dictionary
    mylist = list( dict.fromkeys(abbr_list))
    #sorting in alphabetical order
    mylist.sort()
    #Writting into a file
    document = Document()
    document.add_heading('Appendix - glossary', 0)
    
    for abbr in mylist:
        p = document.add_paragraph(abbr)
    document.save('Appendix-Glossary.docx')
